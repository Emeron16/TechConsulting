from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Colour palette ─────────────────────────────────────────────────────────────
NAVY       = RGBColor(0x0D, 0x1B, 0x3E)
BLUE       = RGBColor(0x00, 0x74, 0xD9)
LIGHTBLUE  = RGBColor(0xE8, 0xF4, 0xFF)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
GREY       = RGBColor(0x64, 0x74, 0x87)
LIGHTGREY  = RGBColor(0xF4, 0xF6, 0xF9)
RED        = RGBColor(0xD6, 0x2B, 0x2B)
GREEN      = RGBColor(0x1A, 0x7A, 0x4A)
AMBER      = RGBColor(0xE6, 0x8A, 0x00)

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

blank = prs.slide_layouts[6]

# ── Helper functions ───────────────────────────────────────────────────────────
def add_rect(slide, l, t, w, h, fill=None, line=None, line_width=Pt(0)):
    shape = slide.shapes.add_shape(1, l, t, w, h)
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, l, t, w, h,
             size=18, bold=False, color=NAVY,
             align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txBox = slide.shapes.add_textbox(l, t, w, h)
    tf    = txBox.text_frame
    tf.word_wrap = wrap
    p  = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return txBox

def add_image(slide, path, l, t, w, h):
    slide.shapes.add_picture(path, l, t, w, h)

def section_header(slide, label):
    add_text(slide, label.upper(),
             Inches(0.5), Inches(0.25), Inches(6), Inches(0.35),
             size=9, bold=True, color=BLUE)

def slide_title(slide, title, subtitle=None):
    add_text(slide, title,
             Inches(0.5), Inches(0.7), Inches(12.3), Inches(0.75),
             size=28, bold=True, color=NAVY)
    add_rect(slide, Inches(0.5), Inches(1.5), Inches(12.33), Inches(0.04), fill=BLUE)
    if subtitle:
        add_text(slide, subtitle,
                 Inches(0.5), Inches(1.6), Inches(12.3), Inches(0.5),
                 size=14, color=GREY)

def bullet_slide(slide, title, bullets, section='', subtitle=None, note=None):
    add_rect(slide, 0, 0, W, H, fill=WHITE)
    if section:
        section_header(slide, section)
    slide_title(slide, title, subtitle)
    top = Inches(2.1)
    for bullet in bullets:
        indent = bullet.startswith('  ')
        txt    = bullet.lstrip()
        bsize  = 14 if not indent else 12
        bcol   = NAVY if not indent else GREY
        prefix = '• ' if not indent else '   – '
        add_text(slide, prefix + txt,
                 Inches(0.6), top, Inches(12.1), Inches(0.45),
                 size=bsize, color=bcol)
        top += Inches(0.42) if not indent else Inches(0.38)
    if note:
        add_rect(slide, Inches(0.5), Inches(6.6), Inches(12.33), Inches(0.65),
                 fill=LIGHTBLUE, line=BLUE, line_width=Pt(1))
        add_text(slide, '⚠  ' + note,
                 Inches(0.65), Inches(6.65), Inches(12.0), Inches(0.55),
                 size=11, color=NAVY, italic=True)

def two_col_slide(slide, title, left_items, right_items,
                  left_head='', right_head='', section=''):
    add_rect(slide, 0, 0, W, H, fill=WHITE)
    if section:
        section_header(slide, section)
    slide_title(slide, title)
    mid = Inches(6.8)

    if left_head:
        add_rect(slide, Inches(0.5), Inches(2.05), Inches(5.9), Inches(0.35), fill=NAVY)
        add_text(slide, left_head, Inches(0.6), Inches(2.07), Inches(5.7), Inches(0.32),
                 size=11, bold=True, color=WHITE)
    if right_head:
        add_rect(slide, mid, Inches(2.05), Inches(5.9), Inches(0.35), fill=NAVY)
        add_text(slide, right_head, Inches(6.9), Inches(2.07), Inches(5.7), Inches(0.32),
                 size=11, bold=True, color=WHITE)

    top_l = top_r = Inches(2.5)
    for item in left_items:
        add_text(slide, '• ' + item, Inches(0.6), top_l, Inches(5.8), Inches(0.45),
                 size=13, color=NAVY)
        top_l += Inches(0.42)
    for item in right_items:
        add_text(slide, '• ' + item, Inches(6.9), top_r, Inches(5.8), Inches(0.45),
                 size=13, color=NAVY)
        top_r += Inches(0.42)

def metric_slide(slide, title, metrics, section=''):
    add_rect(slide, 0, 0, W, H, fill=WHITE)
    if section:
        section_header(slide, section)
    slide_title(slide, title)
    n   = len(metrics)
    box_w = Inches(12.33 / n - 0.2)
    left  = Inches(0.5)
    for label, value, sub, col in metrics:
        add_rect(slide, left, Inches(2.1), box_w, Inches(3.8), fill=LIGHTBLUE, line=col, line_width=Pt(2))
        add_text(slide, label,  left + Inches(0.15), Inches(2.3),  box_w - Inches(0.3), Inches(0.5),
                 size=13, bold=True, color=GREY, align=PP_ALIGN.CENTER)
        add_text(slide, value,  left + Inches(0.15), Inches(3.0),  box_w - Inches(0.3), Inches(1.4),
                 size=36, bold=True, color=col,  align=PP_ALIGN.CENTER)
        add_text(slide, sub,    left + Inches(0.15), Inches(4.5),  box_w - Inches(0.3), Inches(0.8),
                 size=11, color=GREY, align=PP_ALIGN.CENTER)
        left += box_w + Inches(0.2)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Cover
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
add_rect(s, 0, 0, W, H, fill=WHITE)
add_rect(s, 0, 0, Inches(0.18), H, fill=BLUE)
add_rect(s, 0, Inches(6.9), W, Inches(0.6), fill=NAVY)

add_text(s, 'FRAUD TRANSACTION',
         Inches(0.5), Inches(1.6), Inches(10), Inches(1.0),
         size=44, bold=True, color=NAVY)
add_text(s, 'DETECTION',
         Inches(0.5), Inches(2.55), Inches(10), Inches(1.0),
         size=44, bold=True, color=BLUE)
add_text(s, 'A Machine Learning Pipeline for Binary Fraud Classification',
         Inches(0.5), Inches(3.65), Inches(10), Inches(0.6),
         size=16, color=GREY, italic=True)
add_rect(s, Inches(0.5), Inches(4.45), Inches(5), Inches(0.04), fill=LIGHTBLUE)
add_text(s, 'Presenter:  Emeron Marcelle',
         Inches(0.5), Inches(4.6), Inches(8), Inches(0.4),
         size=13, color=NAVY)
add_text(s, 'Date:  June 18, 2026',
         Inches(0.5), Inches(5.0), Inches(8), Inches(0.4),
         size=13, color=GREY)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Agenda
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
bullet_slide(s, 'Agenda', [
    'Business Problem & Dataset Overview',
    'Data Cleaning & Missing Value Strategy',
    'Feature Engineering & Encoding',
    'Exploratory Data Analysis — Key Findings (with charts)',
    'Class Imbalance & Why It Matters',
    'Model Selection & Bayesian Hyperparameter Search',
    'Results: Classification Reports, PR Curve & Feature Importance',
    'Honest Assessment: Perfect Scores & Synthetic Data',
    'Deployed Inference App — Threshold Tuning & SHAP Explanations',
    'Limitations & Next Steps',
], section='Overview')

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Business Problem
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
bullet_slide(s, 'Business Problem', [
    'Financial fraud costs institutions billions annually — detection must be fast and accurate',
    'Challenge: Fraudulent transactions are rare (~8% of all transactions)',
    'Goal: Build a binary classifier to flag fraudulent transactions at inference time',
    'Constraint: Model must be explainable to both technical and business audiences',
    '',
    'Dataset: 15,000 transactions · 18 features · domains: device, location, auth, behavioural',
    '  Target: is_fraudulent_transaction (1 = fraud, 0 = legitimate)',
    '  Class split: 13,800 legitimate vs 1,200 fraud',
], section='Context')

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Dataset Overview
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
two_col_slide(s, 'Dataset Feature Overview',
    left_items=[
        'transaction_amount — USD value',
        'transaction_hour — time of day (0–23)',
        'velocity_last_1hr — txns in last hour',
        'failed_auth_attempts — prior failures',
        'previous_fraud_flags — account history',
        'distance_from_home_km — geo signal',
        'amount_deviation_ratio — vs. avg spend',
        'account_age_days — account maturity',
    ],
    right_items=[
        'device_type — desktop / mobile / tablet',
        'card_type — credit / debit / prepaid / virtual',
        'merchant_category — 8 categories',
        'transaction_channel — online / POS / ATM',
        'location_mismatch — binary flag',
        'ip_country_match — binary flag',
        'cvv_match — binary flag',
        'is_weekend — binary flag',
    ],
    left_head='Numeric Features',
    right_head='Categorical & Binary Features',
    section='Data')

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Missing Data Summary
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
bullet_slide(s, 'Missing Data — Scale & Strategy', [
    'Heavy missingness (~24–25%): transaction_amount, velocity_last_1hr, account_age_days,',
    '  failed_auth_attempts, amount_deviation_ratio, distance_from_home_km, cvv_match',
    '',
    'Moderate missingness (~9–10%): transaction_hour, device_type, location_mismatch,',
    '  previous_fraud_flags, merchant_category, card_type, ip_country_match, is_weekend',
    '',
    'Strategy chosen based on 3 factors: data type · distribution shape · fraud relevance',
    '  MICE — continuous interrelated fraud signals (preserves correlation structure)',
    '  Median — skewed numerics & binary flags (robust to outliers)',
    '  Mode — low-variance binary & categorical columns',
], section='Data Cleaning')

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Why MICE
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
bullet_slide(s, 'Why MICE for Continuous Features?', [
    'MICE = Multiple Imputation by Chained Equations',
    'Models each missing column as a function of all other columns — iterates until stable',
    '',
    'Key insight: fraud-signal features are interrelated:',
    '  High velocity_last_1hr often co-occurs with high failed_auth_attempts',
    '  High previous_fraud_flags correlates with distance_from_home_km',
    '',
    'Simple median imputation would break these correlations — MICE preserves them',
    'Applied to: velocity_last_1hr · account_age_days · failed_auth_attempts ·',
    '  amount_deviation_ratio · previous_fraud_flags · distance_from_home_km',
], section='Data Cleaning')

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Missingness as a Signal
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
bullet_slide(s, 'Missingness as a Fraud Signal', [
    'Key insight: the absence of certain values may itself indicate fraud',
    '',
    'cvv_match missing → CVV was never provided — suspicious behaviour',
    'transaction_amount missing → incomplete transaction record — unusual',
    'ip_country_match missing → IP was masked or unavailable — red flag',
    '',
    'Solution: created binary indicator columns BEFORE imputation:',
    '  cvv_match_missing · transaction_amount_missing · ip_country_match_missing',
    '',
    'These indicators are kept as features — the model can learn from the fact',
    '  that a value was missing, even after the missing value is filled',
], section='Feature Engineering',
note='Indicator columns must always be created before imputation — doing it after results in all zeros.')

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — EDA: Class Distribution + Amount Distribution (charts)
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
add_rect(s, 0, 0, W, H, fill=WHITE)
section_header(s, 'EDA')
slide_title(s, 'Class Imbalance & Transaction Amount Distribution')
add_image(s, 'chart_class_dist.png',  Inches(0.4),  Inches(1.7), Inches(4.6), Inches(3.8))
add_image(s, 'chart_amount_dist.png', Inches(5.2),  Inches(1.7), Inches(7.8), Inches(3.8))
add_rect(s, Inches(0.4), Inches(5.6), Inches(12.33), Inches(1.6), fill=LIGHTBLUE, line=BLUE, line_width=Pt(1))
add_text(s, '• 92% legitimate vs 8% fraud — accuracy alone would be misleading (predicting all-legit = 92%)\n'
            '• transaction_amount is heavily right-skewed: mean $20,898 vs median $126, max $995K\n'
            '• Log-transform normalises the distribution — recommended before production deployment',
         Inches(0.6), Inches(5.65), Inches(12.0), Inches(1.5),
         size=11, color=NAVY)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — EDA: Feature Correlation (chart)
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
add_rect(s, 0, 0, W, H, fill=WHITE)
section_header(s, 'EDA')
slide_title(s, 'Feature Correlation with Fraud Target')
add_image(s, 'chart_correlation.png', Inches(0.5), Inches(1.7), Inches(8.5), Inches(4.8))
add_rect(s, Inches(9.2), Inches(1.9), Inches(3.8), Inches(4.6), fill=LIGHTBLUE, line=BLUE, line_width=Pt(1))
add_text(s, 'Top Fraud Signals',
         Inches(9.3), Inches(2.0), Inches(3.6), Inches(0.4),
         size=12, bold=True, color=NAVY)
bullets = [
    ('velocity_last_1hr',      'r = 0.82'),
    ('failed_auth_attempts',   'r = 0.81'),
    ('previous_fraud_flags',   'r = 0.78'),
    ('distance_from_home_km',  'r = 0.73'),
    ('ip_country_match',       'r = 0.64'),
    ('location_mismatch',      'r = 0.59'),
]
top = Inches(2.5)
for feat, corr in bullets:
    add_text(s, f'• {feat}',  Inches(9.3), top, Inches(2.6), Inches(0.35), size=10, color=NAVY)
    add_text(s, corr,         Inches(12.0), top, Inches(0.9), Inches(0.35), size=10, bold=True, color=RED)
    top += Inches(0.38)
add_text(s, 'Behavioural signals dominate.\nNo multicollinearity detected\n(all inter-feature r < 0.8)',
         Inches(9.3), Inches(5.0), Inches(3.6), Inches(0.9),
         size=10, color=GREY, italic=True)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — EDA: Fraud vs Legit violin plots (chart)
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
add_rect(s, 0, 0, W, H, fill=WHITE)
section_header(s, 'EDA')
slide_title(s, 'Top Fraud Signals — Distribution by Class')
add_image(s, 'chart_violin.png', Inches(0.5), Inches(1.7), Inches(12.33), Inches(4.5))
add_rect(s, Inches(0.5), Inches(6.3), Inches(12.33), Inches(0.9), fill=LIGHTBLUE, line=BLUE, line_width=Pt(1))
add_text(s, '• Fraud transactions cluster at HIGH values for all four signals — clean separation in synthetic data\n'
            '• Legitimate transactions stay near zero — confirms these are the most discriminative features',
         Inches(0.65), Inches(6.35), Inches(12.0), Inches(0.8),
         size=11, color=NAVY)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Class Imbalance
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
bullet_slide(s, 'Class Imbalance — Why It Matters', [
    'Dataset: 13,800 legitimate (92%) vs 1,200 fraud (8%)',
    '',
    'A naive model predicting "not fraud" every time achieves 92% accuracy',
    '  — but catches zero fraud cases. Accuracy is the wrong metric here.',
    '',
    'Solution: SMOTE (Synthetic Minority Oversampling Technique)',
    '  Generates synthetic fraud samples by interpolating between real fraud cases',
    '  After SMOTE: 11,040 legitimate vs 11,040 fraud in training set',
    '',
    'Critical rule: SMOTE applied to training set ONLY — never before splitting',
    '  SMOTE placed inside the pipeline so it runs fresh per CV fold',
    '  Prevents synthetic samples from leaking into validation or test sets',
], section='Modelling',
note='Metrics used: Precision · Recall · F1 · PR-AUC — not accuracy.')

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — Model Selection
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
two_col_slide(s, 'Model Selection — Simple to Complex',
    left_items=[
        'Logistic Regression — baseline',
        'Linear decision boundary',
        'Fully interpretable coefficients',
        'StandardScaler for feature normalisation',
        'Fast — establishes performance floor',
        '',
        'Random Forest — robust middle ground',
        'Ensemble of decision trees',
        'Handles outliers & mixed features',
    ],
    right_items=[
        'XGBoost — high-performance benchmark',
        'Sequential error correction',
        'scale_pos_weight for imbalance',
        'Best performer on tabular fraud data',
        '',
        'Bayesian Hyperparameter Search',
        'Smarter than Grid / Random search',
        'Learns from each trial — 30 iterations',
        'SMOTE baked into each pipeline',
    ],
    left_head='Algorithms Tested',
    right_head='Optimisation Strategy',
    section='Modelling')

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — Results metrics table
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
add_rect(s, 0, 0, W, H, fill=WHITE)
section_header(s, 'Results')
slide_title(s, 'Model Performance — All Three Models')

headers = ['Model', 'CV AUC', 'ROC-AUC', 'PR-AUC', 'F1 (Fraud)']
rows    = [
    ['Logistic Regression', '1.0000', '1.0000', '1.0000', '1.00'],
    ['Random Forest',       '1.0000', '1.0000', '1.0000', '1.00'],
    ['XGBoost',             '1.0000', '1.0000', '1.0000', '1.00'],
]
col_w = [Inches(3.2), Inches(2.1), Inches(2.1), Inches(2.1), Inches(2.1)]
lefts = [Inches(0.4)]
for w in col_w[:-1]:
    lefts.append(lefts[-1] + w + Inches(0.06))

top = Inches(2.1)
for i, h in enumerate(headers):
    add_rect(s, lefts[i], top, col_w[i], Inches(0.38), fill=NAVY)
    add_text(s, h, lefts[i] + Inches(0.05), top + Inches(0.04),
             col_w[i] - Inches(0.1), Inches(0.32),
             size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

for ri, row in enumerate(rows):
    bg = LIGHTBLUE if ri % 2 == 0 else WHITE
    top_r = top + Inches(0.42) * (ri + 1)
    for ci, cell in enumerate(row):
        add_rect(s, lefts[ci], top_r, col_w[ci], Inches(0.38), fill=bg, line=LIGHTGREY, line_width=Pt(0.5))
        add_text(s, cell, lefts[ci] + Inches(0.05), top_r + Inches(0.04),
                 col_w[ci] - Inches(0.1), Inches(0.32),
                 size=12, color=NAVY, align=PP_ALIGN.CENTER)

add_rect(s, Inches(0.4), Inches(4.0), Inches(12.33), Inches(1.2),
         fill=RGBColor(0xFF,0xF3,0xCD), line=AMBER, line_width=Pt(1.5))
add_text(s, '⚠  Perfect scores (1.0000) across all models warrant scrutiny. '
            'This dataset is synthetic — fraud patterns are too clean and consistent '
            'compared to real-world data. The model architecture and pipeline are sound; '
            'results should be validated on a real-world dataset before production deployment.',
         Inches(0.6), Inches(4.05), Inches(12.0), Inches(1.1),
         size=11, color=NAVY, italic=True)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — PR Curve (chart)
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
add_rect(s, 0, 0, W, H, fill=WHITE)
section_header(s, 'Results')
slide_title(s, 'Precision-Recall Curve & Threshold Tuning')
add_image(s, 'chart_pr_curve.png', Inches(0.5), Inches(1.7), Inches(7.5), Inches(4.8))
add_rect(s, Inches(8.2), Inches(1.9), Inches(4.8), Inches(4.6), fill=LIGHTBLUE, line=BLUE, line_width=Pt(1))
add_text(s, 'Why Tune the Threshold?',
         Inches(8.35), Inches(2.05), Inches(4.5), Inches(0.4),
         size=12, bold=True, color=NAVY)
pts = [
    'Default threshold = 0.5 is arbitrary',
    'Optimal threshold maximises F1:',
    '   balances precision vs recall',
    'Tuned threshold = 0.41',
    '  → catches more fraud (higher recall)',
    '  → accepts slightly more false positives',
    '',
    'In fraud detection, missing a fraud case',
    'is more costly than a false alarm',
    '→ tuning toward recall is deliberate',
]
top = Inches(2.55)
for pt in pts:
    sz = 11 if not pt.startswith('Tuned') else 12
    bold = pt.startswith('Tuned')
    col = RED if pt.startswith('Tuned') else NAVY
    add_text(s, pt, Inches(8.35), top, Inches(4.5), Inches(0.38),
             size=sz, color=col, bold=bold)
    top += Inches(0.35)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 15 — Honest Assessment
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
bullet_slide(s, 'Honest Assessment — Perfect Scores', [
    'All three models achieved 1.0 AUC, F1, Precision and Recall on the test set',
    '',
    'This is NOT data leakage — confirmed by single-feature AUC analysis:',
    '  Leakage would mean a feature encodes the answer (e.g. fraud confirmed after the fact)',
    '  Our top features (velocity, failed_auth, previous_flags) are available at inference time',
    '',
    'Root cause: the dataset is synthetic',
    '  Synthetic fraud datasets are designed with clean, consistent fraud patterns',
    '  Real fraudsters adapt — legitimate users sometimes behave suspiciously',
    '  Real datasets produce AUC in the 0.85–0.95 range, not 1.0',
    '',
    'Conclusion: pipeline and methodology are correct — on realistic data,',
    '  the performance will reflect true model capability',
], section='Results',
note='Recommendation: test on a real financial fraud dataset (e.g. IEEE-CIS or PaySim) before production.')

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 16 — Feature Importance
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
bullet_slide(s, 'Feature Importance — Key Fraud Signals', [
    'Top predictors consistent across all three models:',
    '',
    '  1.  velocity_last_1hr         r = 0.82 — rapid transactions = fraud signal',
    '  2.  failed_auth_attempts      r = 0.81 — repeated authentication failures',
    '  3.  previous_fraud_flags      r = 0.78 — account history is highly predictive',
    '  4.  distance_from_home_km     r = 0.73 — geographic anomaly detection',
    '  5.  ip_country_match          r = 0.64 — IP origin mismatch',
    '  6.  location_mismatch         r = 0.59 — registered address vs transaction location',
    '  7.  cvv_match                 r = 0.57 — failed CVV verification',
    '',
    'Business insight: behavioural signals (velocity, auth failures) outperform',
    '  static account signals (account age, card type) for fraud detection',
], section='Results')

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 17 — Why Logistic Regression
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
bullet_slide(s, 'Why Logistic Regression for Deployment?', [
    'All three models performed identically on this synthetic dataset',
    'Logistic Regression chosen for deployment for three reasons:',
    '',
    '  Interpretability: coefficients directly map features to fraud probability',
    '    Business stakeholders can understand and audit the decision logic',
    '',
    '  Speed: inference in microseconds — critical for real-time transaction scoring',
    '    XGBoost and Random Forest are 10–100x slower at inference time',
    '',
    '  Regulatory compliance: explainable models preferred in financial services',
    '    Easier to justify decisions to regulators (e.g. GDPR, SR 11-7)',
    '',
    'On a real-world dataset with more noise, XGBoost would likely outperform LR',
    '  and should be re-evaluated before production',
], section='Deployment')

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 18 — Streamlit App (updated)
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
bullet_slide(s, 'Deployed Inference App — Streamlit', [
    'Real-time inference web app — user enters 16 raw transaction fields via form',
    'Loads saved LR pipeline (fraud_lr_model.pkl) + tuned threshold (model_threshold.pkl)',
    '',
    'App applies the same preprocessing pipeline as the notebook:',
    '  Creates indicator columns → one-hot encodes → aligns to 31-column model input',
    '  StandardScaler applied before prediction (same scaler used in training)',
    '',
    'Output 1 — Risk Assessment:',
    '  Fraud probability % · Decision threshold (0.41) · LEGITIMATE / FRAUD DETECTED verdict',
    '',
    'Output 2 — SHAP-style Explanation:',
    '  Top 10 feature contributions (coef × scaled value) shown as colour-coded bar chart',
    '  Red bars pushed toward fraud · Green bars pushed toward legitimate',
    '  Key Drivers section lists top 3 features in each direction with input values',
], section='Deployment')

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 19 — Pipeline Summary
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
add_rect(s, 0, 0, W, H, fill=WHITE)
section_header(s, 'Summary')
slide_title(s, 'End-to-End Pipeline')

steps = [
    ('1', 'Raw Data',        '15K rows · 18 features · 8% fraud'),
    ('2', 'Cleaning',        'Drop ID · MICE · Median · Mode imputation'),
    ('3', 'Feature Eng.',    'Missingness indicators · One-hot encoding'),
    ('4', 'EDA',             'Skewness · Correlation · Class imbalance charts'),
    ('5', 'Modelling',       'SMOTE in pipeline · Bayesian search · 3 models'),
    ('6', 'Evaluation',      'Classification report · PR curve · Threshold tuning'),
    ('7', 'Deployment',      'LR + Scaler · Streamlit app · SHAP explanations'),
]

box_w = Inches(1.55)
gap   = Inches(0.18)
left  = Inches(0.35)
top_b = Inches(2.2)

for i, (num, title, sub) in enumerate(steps):
    add_rect(s, left, top_b, box_w, Inches(3.2), fill=NAVY if i == 0 else LIGHTBLUE,
             line=BLUE, line_width=Pt(1.5))
    add_text(s, num,   left + Inches(0.05), top_b + Inches(0.1),
             box_w - Inches(0.1), Inches(0.45),
             size=22, bold=True, color=BLUE if i > 0 else WHITE, align=PP_ALIGN.CENTER)
    add_text(s, title, left + Inches(0.05), top_b + Inches(0.6),
             box_w - Inches(0.1), Inches(0.5),
             size=12, bold=True, color=WHITE if i == 0 else NAVY, align=PP_ALIGN.CENTER)
    add_text(s, sub,   left + Inches(0.05), top_b + Inches(1.2),
             box_w - Inches(0.1), Inches(1.8),
             size=9, color=WHITE if i == 0 else GREY, align=PP_ALIGN.CENTER)
    if i < len(steps) - 1:
        add_text(s, '▶', left + box_w + Inches(0.02), top_b + Inches(1.3),
                 gap, Inches(0.4), size=12, color=BLUE, align=PP_ALIGN.CENTER)
    left += box_w + gap

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 20 — Limitations & Next Steps (updated)
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
add_rect(s, 0, 0, W, H, fill=WHITE)
section_header(s, 'Next Steps')
slide_title(s, 'Limitations & Next Steps')
mid = Inches(6.8)

# LEFT — Remaining Limitations
add_rect(s, Inches(0.5), Inches(2.05), Inches(5.9), Inches(0.35), fill=NAVY)
add_text(s, 'Remaining Limitations', Inches(0.6), Inches(2.07), Inches(5.7), Inches(0.32),
         size=11, bold=True, color=WHITE)

limitations = [
    'Synthetic dataset — fraud patterns are too clean',
    'Perfect scores not representative of production',
    'Single model saved — no A/B infrastructure',
    'No model drift monitoring in production',
    'SHAP approximated via coef × value (not full SHAP)',
]
top = Inches(2.5)
for item in limitations:
    add_text(s, '• ' + item, Inches(0.6), top, Inches(5.8), Inches(0.45), size=13, color=NAVY)
    top += Inches(0.44)

# RESOLVED box
add_rect(s, Inches(0.5), Inches(5.1), Inches(5.9), Inches(1.8),
         fill=RGBColor(0xE6, 0xF4, 0xEA), line=GREEN, line_width=Pt(1.5))
add_text(s, '✓  Already Resolved in This Project',
         Inches(0.65), Inches(5.18), Inches(5.6), Inches(0.35),
         size=11, bold=True, color=GREEN)
resolved = ['StandardScaler added to LR pipeline',
            'Classification threshold tuned via PR curve (0.41)',
            'Feature contributions (SHAP-style) shown in Streamlit app']
top = Inches(5.55)
for r in resolved:
    add_text(s, '✓  ' + r, Inches(0.65), top, Inches(5.6), Inches(0.35), size=11, color=GREEN)
    top += Inches(0.36)

# RIGHT — Recommended Next Steps
add_rect(s, mid, Inches(2.05), Inches(5.9), Inches(0.35), fill=NAVY)
add_text(s, 'Recommended Next Steps', Inches(6.9), Inches(2.07), Inches(5.7), Inches(0.32),
         size=11, bold=True, color=WHITE)

nextsteps = [
    'Validate on real dataset (IEEE-CIS / PaySim)',
    'Add LightGBM to model comparison',
    'Set up model monitoring for drift detection',
    'Build proper A/B model testing infrastructure',
    'Add full SHAP background data for global explanations',
    'Containerise app with Docker for cloud deployment',
    'Add authentication layer to Streamlit app',
]
top = Inches(2.5)
for item in nextsteps:
    add_text(s, '• ' + item, Inches(6.9), top, Inches(5.8), Inches(0.45), size=13, color=NAVY)
    top += Inches(0.44)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 21 — Conclusions
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
add_rect(s, 0, 0, W, H, fill=WHITE)
add_rect(s, 0, 0, Inches(0.18), H, fill=BLUE)
section_header(s, 'Conclusions')
slide_title(s, 'Key Takeaways')

takeaways = [
    ('Pipeline',       'Full end-to-end ML pipeline built: cleaning → imputation → encoding → modelling → deployment'),
    ('Imbalance',      'SMOTE inside the pipeline correctly handles 8% fraud rate without data leakage'),
    ('Signals',        'velocity_last_1hr, failed_auth_attempts, and previous_fraud_flags are the strongest fraud predictors'),
    ('Improvements',   'StandardScaler, threshold tuning (0.41) and SHAP-style explanations all implemented and deployed'),
    ('Deployment',     'LR pipeline deployed as a real-time Streamlit app with probability, threshold, and feature explanations'),
    ('Honesty',        'Perfect 1.0 scores reflect a synthetic dataset — methodology is sound, not the scores'),
]

top = Inches(2.0)
for label, text in takeaways:
    add_rect(s, Inches(0.5), top, Inches(1.8), Inches(0.55), fill=BLUE)
    add_text(s, label, Inches(0.5), top + Inches(0.06), Inches(1.8), Inches(0.45),
             size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, text, Inches(2.45), top + Inches(0.06), Inches(10.2), Inches(0.45),
             size=12, color=NAVY)
    top += Inches(0.68)

# ══════════════════════════════════════════════════════════════════════════════
# Save
# ══════════════════════════════════════════════════════════════════════════════
out = 'Fraud_Detection_Presentation.pptx'
prs.save(out)
print(f"Saved: {out}  ({len(prs.slides)} slides)")
